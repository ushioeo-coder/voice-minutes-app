import streamlit as st
from google import genai
from google.genai import types
import tempfile
import os
import time
from datetime import datetime
import json
import requests
import hashlib
import shutil
import zipfile
import io
from PIL import Image, ImageDraw, ImageFont

# pydub & ffmpeg (非同期でimport)
try:
    from pydub import AudioSegment
    import imageio_ffmpeg
    
    # FFmpegのパスをimageio_ffmpegから取得して設定
    # これによりシステムにffmpegがなくても動作する
    AudioSegment.converter = imageio_ffmpeg.get_ffmpeg_exe()
    PYDUB_AVAILABLE = True
except ImportError:
    PYDUB_AVAILABLE = False

# Firebase初期化（オプション）
try:
    import firebase_admin
    from firebase_admin import credentials, firestore
    FIREBASE_ADMIN_AVAILABLE = True
except ImportError:
    FIREBASE_ADMIN_AVAILABLE = False

# ページ設定
st.set_page_config(
    page_title="MinuteSlide",
    page_icon="📊",
    layout="wide"
)

# システムプロンプト
SYSTEM_PROMPT = """あなたはプロの議事録作成者です。提供された音声ファイルを注意深く聴き、以下の構成で議事録をまとめてください。

## 会議の目的
何のための会議か。

## 議論の要約
どのような議論が行われたか（時系列またはトピック別）。

## ネクストアクション
- [ ] 【担当者名】具体的なタスク内容（あれば期限も）

## 決定事項
合意に至った内容。

---
**注意点:**
- 発話者が特定できる場合は「Aさん：〜」のように記載してください。
- 音質が悪い箇所は推測せず、事実に忠実であること。
- 不明瞭な部分は「（聞き取り不可）」と記載してください。
"""

# スライド生成プロンプト
SLIDE_PROMPT = """以下の議事録をプレゼンテーション用のスライドに変換してください。

出力形式は以下のJSON構造で返してください（```json と ``` で囲んでください）：

```json
{
  "title": "会議のタイトル",
  "date": "会議日付（議事録から推測）",
  "slides": [
    {
      "type": "title",
      "title": "メインタイトル",
      "subtitle": "サブタイトル（日付など）"
    },
    {
      "type": "purpose",
      "title": "会議の目的",
      "content": "目的の説明"
    },
    {
      "type": "discussion",
      "title": "議論トピック名",
      "bullets": ["ポイント1", "ポイント2", "ポイント3"]
    },
    {
      "type": "actions",
      "title": "ネクストアクション",
      "items": [
        {"assignee": "担当者", "task": "タスク内容", "deadline": "期限（あれば）"}
      ]
    },
    {
      "type": "decisions",
      "title": "決定事項",
      "bullets": ["決定1", "決定2"]
    }
  ]
}
```

注意点:
- 各スライドは1つのトピックに絞り、簡潔に
- 箇条書きは最大5項目まで
- 議論が複数トピックある場合は複数のdiscussionスライドを作成
- JSONのみを出力してください
"""

# Firebase Auth REST API エンドポイント
FIREBASE_AUTH_URL = "https://identitytoolkit.googleapis.com/v1/accounts"

def generate_org_id(org_name: str) -> str:
    """組織名から一意のIDを生成"""
    return hashlib.md5(org_name.lower().strip().encode()).hexdigest()[:12]

def firebase_auth_request(endpoint: str, data: dict, api_key: str):
    """Firebase Auth REST APIリクエスト（リトライ付き）"""
    url = f"{FIREBASE_AUTH_URL}:{endpoint}?key={api_key}"
    
    # セッションとリトライの設定
    session = requests.Session()
    adapter = requests.adapters.HTTPAdapter(max_retries=3)
    session.mount('https://', adapter)
    
    try:
        response = session.post(url, json=data, timeout=10)
        return response.json()
    except requests.exceptions.RequestException as e:
        st.error(f"ネットワークエラーが発生しました: {e}")
        return {"error": {"message": "Network Error"}}

def sign_up(email: str, password: str, api_key: str):
    """新規ユーザー登録"""
    data = {
        "email": email,
        "password": password,
        "returnSecureToken": True
    }
    return firebase_auth_request("signUp", data, api_key)

def sign_in(email: str, password: str, api_key: str):
    """ログイン"""
    data = {
        "email": email,
        "password": password,
        "returnSecureToken": True
    }
    return firebase_auth_request("signInWithPassword", data, api_key)

def reset_password(email: str, api_key: str):
    """パスワードリセットメール送信"""
    data = {
        "email": email,
        "requestType": "PASSWORD_RESET"
    }
    return firebase_auth_request("sendOobCode", data, api_key)

def init_firestore():
    """Firestore初期化"""
    if not FIREBASE_ADMIN_AVAILABLE:
        return None
    
    if firebase_admin._apps:
        return firestore.client()
    
    try:
        if "firebase" in st.secrets and "service_account" in st.secrets["firebase"]:
            cred_dict = json.loads(st.secrets["firebase"]["service_account"])
            cred = credentials.Certificate(cred_dict)
            firebase_admin.initialize_app(cred)
            return firestore.client()
    except Exception as e:
        st.warning(f"Firestore未設定: {e}")
    return None

def get_or_create_organization(db, org_name: str):
    """組織を取得または作成"""
    if db is None:
        return None
    
    org_id = generate_org_id(org_name)
    org_ref = db.collection("organizations").document(org_id)
    org_doc = org_ref.get()
    
    if not org_doc.exists:
        # 新規組織作成
        org_ref.set({
            "name": org_name.strip(),
            "created_at": datetime.now()
        })
    
    return org_id

def get_user_organization(db, user_id: str):
    """ユーザーの組織IDを取得"""
    if db is None:
        return None
    
    try:
        user_doc = db.collection("users").document(user_id).get()
        if user_doc.exists:
            return user_doc.to_dict().get("organization_id")
    except:
        pass
    return None

def save_user_organization(db, user_id: str, user_email: str, org_id: str, org_name: str):
    """ユーザーの組織情報を保存"""
    if db is None:
        return
    
    try:
        db.collection("users").document(user_id).set({
            "email": user_email,
            "organization_id": org_id,
            "organization_name": org_name,
            "created_at": datetime.now()
        })
    except Exception as e:
        st.error(f"ユーザー情報保存エラー: {e}")

def get_organization_name(db, org_id: str):
    """組織IDから組織名を取得"""
    if db is None:
        return None
    
    try:
        org_doc = db.collection("organizations").document(org_id).get()
        if org_doc.exists:
            return org_doc.to_dict().get("name")
    except:
        pass
    return None

def check_firebase_auth(db):
    """Firebase Authentication認証（組織対応）"""
    # セッション状態初期化
    if "authenticated" not in st.session_state:
        st.session_state.authenticated = False
        st.session_state.user_email = None
        st.session_state.user_id = None
        st.session_state.organization_id = None
        st.session_state.organization_name = None
    
    # Firebase API Key確認
    if "firebase" not in st.secrets or "api_key" not in st.secrets["firebase"]:
        st.error("Firebase API Keyが設定されていません")
        st.info("secrets.tomlに[firebase]セクションを追加してください")
        return False
    
    api_key = st.secrets["firebase"]["api_key"]
    
    if st.session_state.authenticated:
        # 組織情報がまだない場合は取得
        if st.session_state.organization_id is None and db:
            org_id = get_user_organization(db, st.session_state.user_id)
            if org_id:
                st.session_state.organization_id = org_id
                st.session_state.organization_name = get_organization_name(db, org_id)
        return True
    
    # ログイン/新規登録タブ
    tab1, tab2, tab3 = st.tabs(["🔑 ログイン", "📝 新規登録", "🔄 パスワードリセット"])
    
    with tab1:
        st.subheader("ログイン")
        email = st.text_input("メールアドレス", key="login_email")
        password = st.text_input("パスワード", type="password", key="login_password")
        
        if st.button("ログイン", type="primary", use_container_width=True):
            if email and password:
                with st.spinner("ログイン中..."):
                    result = sign_in(email, password, api_key)
                
                if "idToken" in result:
                    user_id = result.get("localId")
                    st.session_state.authenticated = True
                    st.session_state.user_email = result.get("email")
                    st.session_state.user_id = user_id
                    
                    # 組織情報を取得
                    if db:
                        org_id = get_user_organization(db, user_id)
                        if org_id:
                            st.session_state.organization_id = org_id
                            st.session_state.organization_name = get_organization_name(db, org_id)
                        else:
                            st.warning("組織情報がありません。新規登録し直すか、管理者にお問い合わせください。")
                    
                    st.success("ログイン成功！")
                    st.rerun()
                else:
                    error_msg = result.get("error", {}).get("message", "不明なエラー")
                    if "INVALID_LOGIN_CREDENTIALS" in error_msg:
                        st.error("メールアドレスまたはパスワードが間違っています")
                    elif "USER_NOT_FOUND" in error_msg:
                        st.error("ユーザーが見つかりません")
                    else:
                        st.error(f"ログインエラー: {error_msg}")
            else:
                st.warning("メールアドレスとパスワードを入力してください")
    
    with tab2:
        st.subheader("新規アカウント登録")
        
        # 組織名入力
        org_name = st.text_input("組織名（会社名・チーム名）", key="signup_org",
                                  help="同じ組織名のメンバーと議事録を共有できます")
        
        new_email = st.text_input("メールアドレス", key="signup_email")
        new_password = st.text_input("パスワード", type="password", key="signup_password", 
                                      help="6文字以上")
        confirm_password = st.text_input("パスワード（確認）", type="password", key="confirm_password")
        
        if st.button("アカウント作成", type="primary", use_container_width=True):
            if org_name and new_email and new_password and confirm_password:
                if new_password != confirm_password:
                    st.error("パスワードが一致しません")
                elif len(new_password) < 6:
                    st.error("パスワードは6文字以上にしてください")
                elif len(org_name.strip()) < 2:
                    st.error("組織名を2文字以上で入力してください")
                else:
                    with st.spinner("アカウント作成中..."):
                        result = sign_up(new_email, new_password, api_key)
                    
                    if "idToken" in result:
                        user_id = result.get("localId")
                        user_email = result.get("email")
                        
                        # 組織を作成または取得
                        if db:
                            org_id = get_or_create_organization(db, org_name)
                            save_user_organization(db, user_id, user_email, org_id, org_name.strip())
                            st.session_state.organization_id = org_id
                            st.session_state.organization_name = org_name.strip()
                        
                        st.session_state.authenticated = True
                        st.session_state.user_email = user_email
                        st.session_state.user_id = user_id
                        st.success("アカウント作成成功！")
                        st.rerun()
                    else:
                        error_msg = result.get("error", {}).get("message", "不明なエラー")
                        if "EMAIL_EXISTS" in error_msg:
                            st.error("このメールアドレスは既に登録されています")
                        elif "WEAK_PASSWORD" in error_msg:
                            st.error("パスワードが弱すぎます。6文字以上にしてください")
                        elif "INVALID_EMAIL" in error_msg:
                            st.error("メールアドレスの形式が正しくありません")
                        else:
                            st.error(f"登録エラー: {error_msg}")
            else:
                st.warning("すべての項目を入力してください")
    
    with tab3:
        st.subheader("パスワードリセット")
        reset_email = st.text_input("登録済みメールアドレス", key="reset_email")
        
        if st.button("リセットメールを送信", use_container_width=True):
            if reset_email:
                with st.spinner("送信中..."):
                    result = reset_password(reset_email, api_key)
                
                if "error" not in result:
                    st.success("パスワードリセットメールを送信しました。メールを確認してください。")
                else:
                    error_msg = result.get("error", {}).get("message", "不明なエラー")
                    if "EMAIL_NOT_FOUND" in error_msg:
                        st.error("このメールアドレスは登録されていません")
                    else:
                        st.error(f"エラー: {error_msg}")
            else:
                st.warning("メールアドレスを入力してください")
    
    return False

def save_to_firestore(db, title, content, audio_filename, user_email, user_id, org_id):
    """議事録をFirestoreに保存（組織別）"""
    if db is None or org_id is None:
        return None
    
    try:
        doc_ref = db.collection("organizations").document(org_id).collection("minutes").document()
        doc_ref.set({
            "title": title,
            "content": content,
            "audio_filename": audio_filename,
            "created_by": user_email,
            "user_id": user_id,
            "created_at": datetime.now()
        })
        return doc_ref.id
        return doc_ref.id
    except Exception as e:
        st.error(f"保存エラー: {e}")
        return None

def delete_minute(db, org_id, minute_id):
    """議事録を削除"""
    if db is None or org_id is None:
        return False
    
    try:
        db.collection("organizations").document(org_id).collection("minutes").document(minute_id).delete()
        return True
    except Exception as e:
        st.error(f"削除エラー: {e}")
        return False

def get_minutes_history(db, org_id):
    """議事録履歴を取得（同一組織のみ）"""
    if db is None or org_id is None:
        return []
    
    try:
        docs = db.collection("organizations").document(org_id).collection("minutes").order_by("created_at", direction=firestore.Query.DESCENDING).limit(20).stream()
        history = []
        for doc in docs:
            data = doc.to_dict()
            data["id"] = doc.id
            history.append(data)
        return history
    except Exception as e:
        st.warning(f"履歴取得エラー: {e}")
        return []

def get_minute_by_id(db, org_id, doc_id):
    """IDで議事録を取得（組織別）"""
    if db is None or org_id is None:
        return None
    
    try:
        doc = db.collection("organizations").document(org_id).collection("minutes").document(doc_id).get()
        if doc.exists:
            data = doc.to_dict()
            data["id"] = doc.id
            return data
    except Exception as e:
        st.error(f"取得エラー: {e}")
    return None

def upload_and_process_audio(audio_file, model_name, additional_instructions):
    """音声ファイルをアップロードしてGeminiで処理"""
    
    # クライアント設定
    client = genai.Client(api_key=st.secrets["api"]["google_api_key"])
    
    # 一時ファイルに保存
    suffix = os.path.splitext(audio_file.name)[1]
    with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp_file:
        tmp_file.write(audio_file.getvalue())
        tmp_path = tmp_file.name
    
    try:
        # 進捗表示
        progress_bar = st.progress(0)
        status_text = st.empty()
        
        # 音声分割処理 (pydubが必要)
        chunk_files = []
        full_transcript = ""
        
        if PYDUB_AVAILABLE:
            try:
                # ファイルポインタをリセット
                audio_file.seek(0)
                audio = AudioSegment.from_file(audio_file)
                duration_ms = len(audio)
                CHUNK_LENGTH_MS = 10 * 60 * 1000  # 10分
                
                if duration_ms > CHUNK_LENGTH_MS:
                    status_text.text(f"⚠️ 長時間音声のため分割処理を行います（約{duration_ms/1000/60:.1f}分）...")
                    chunks_count = (duration_ms // CHUNK_LENGTH_MS) + 1
                    
                    for i in range(chunks_count):
                        start_ms = i * CHUNK_LENGTH_MS
                        end_ms = min((i + 1) * CHUNK_LENGTH_MS, duration_ms)
                        chunk = audio[start_ms:end_ms]
                        
                        # 一時ファイルに保存
                        chunk_name = f"{tmp_path}_part{i}{suffix}"
                        chunk.export(chunk_name, format=suffix.replace('.', ''))
                        chunk_files.append(chunk_name)
                    
                    status_text.text(f"✂️ 音声を {len(chunk_files)} 個のパートに分割しました。順次処理します。")
                    time.sleep(1)
            except Exception as e:
                # 分割処理失敗の詳細を表示
                st.warning(f"⚠️ 音声分割処理に失敗しました（通常モードで続行します）: {e}")
                # st.error(traceback.format_exc()) # 必要なら詳細ログ
                pass
        
        # 分割ファイルがある場合はループ処理
        if chunk_files:
            total_chunks = len(chunk_files)
            
            for i, chunk_path in enumerate(chunk_files):
                status_text.text(f"🔄 パート {i+1}/{total_chunks} を処理中...")
                progress_step = 100 / total_chunks
                current_progress = int(i * progress_step)
                progress_bar.progress(current_progress)
                
                # アップロード
                uploaded_file = client.files.upload(file=chunk_path)
                
                # 処理待機
                while uploaded_file.state.name == "PROCESSING":
                    time.sleep(2)
                    uploaded_file = client.files.get(name=uploaded_file.name)
                
                if uploaded_file.state.name == "FAILED":
                    st.error(f"パート {i+1} の処理に失敗しました。")
                    continue
                
                # 生成 (リトライ付き)
                # ... (既存のリトライロジックをここでも使うため、関数化すべきだが今回はインライン展開)
                chunk_prompt = f"以下の音声は会議の一部（パート {i+1}/{total_chunks}）です。内容を詳細に書き起こして要約してください。\n\n" + SYSTEM_PROMPT
                
                # リトライ処理 (簡易版)
                for attempt in range(3):
                    try:
                        response = client.models.generate_content(
                            model=model_name,
                            contents=[types.Content(role="user", parts=[
                                types.Part.from_uri(file_uri=uploaded_file.uri, mime_type=uploaded_file.mime_type),
                                types.Part.from_text(text=chunk_prompt)
                            ])]
                        )
                        full_transcript += f"\n\n--- Part {i+1} ---\n\n" + response.text
                        break
                    except Exception as e:
                        if "429" in str(e) and attempt < 2:
                            time.sleep(60)
                            continue
                        elif attempt == 2:
                            full_transcript += f"\n\n[Part {i+1} Error: {e}]\n"
                
                # 削除
                try:
                    client.files.delete(name=uploaded_file.name)
                except:
                    pass
                
                # クリーンアップ
                if os.path.exists(chunk_path):
                    os.remove(chunk_path)
            
            # 結合後の再要約リクエスト（オプションだが今回は単純結合の結果を返す）
            progress_bar.progress(100)
            status_text.text("✅ 分割処理完了！")
            return full_transcript

        # --- 以下、通常処理 (分割なし) ---
        
        # 音声ファイルをアップロード
        status_text.text("🔄 音声ファイルをアップロード中...")
        progress_bar.progress(20)
        
        uploaded_file = client.files.upload(file=tmp_path)
        
        # ファイルがアクティブになるまで待機
        status_text.text("⏳ 音声ファイルを処理中...")
        progress_bar.progress(40)
        
        while uploaded_file.state.name == "PROCESSING":
            time.sleep(2)
            uploaded_file = client.files.get(name=uploaded_file.name)
        
        if uploaded_file.state.name == "FAILED":
            st.error("音声ファイルの処理に失敗しました。")
            return None
        
        # 議事録生成
        status_text.text("📝 議事録を生成中...")
        progress_bar.progress(60)
        
        # プロンプト構築
        prompt = SYSTEM_PROMPT
        if additional_instructions:
            prompt += f"\n\n**追加指示:**\n{additional_instructions}"
        
        # リトライ設定（強化版）
        max_retries = 5
        base_delay = 10
        
        for attempt in range(max_retries):
            try:
                response = client.models.generate_content(
                    model=model_name,
                    contents=[
                        types.Content(
                            role="user",
                            parts=[
                                types.Part.from_uri(
                                    file_uri=uploaded_file.uri,
                                    mime_type=uploaded_file.mime_type
                                ),
                                types.Part.from_text(text=prompt)
                            ]
                        )
                    ]
                )
                break  # 成功したらループを抜ける
            except Exception as e:
                if "RESOURCE_EXHAUSTED" in str(e) or "429" in str(e):
                    if attempt < max_retries - 1:
                        # トークン制限（1分間）の回復を待つため、一律60秒待機
                        wait_time = 60
                        status_text.text(f"⚠️ API制限（429）。トークン回復まで{wait_time}秒待機します... ({attempt + 1}/{max_retries})")
                        time.sleep(wait_time)
                        continue
                
                # リトライしない、またはリトライ上限到達
                if "RESOURCE_EXHAUSTED" in str(e) or "429" in str(e):
                    st.error("⚠️ APIの利用制限（アクセス集中）が発生しました。数分時間を置いてから再度お試しいただくか、モデルを「**Gemini 2.0 Flash Lite**」に変更してみてください。")
                    st.error(f"詳細エラー: {e}")
                else:
                    st.error(f"議事録生成エラー: {e}")
                
                # アップロードしたファイルを削除して終了
                try:
                    client.files.delete(name=uploaded_file.name)
                except:
                    pass
                return None
        
        progress_bar.progress(100)
        status_text.text("✅ 完了！")
        
        # アップロードしたファイルを削除
        try:
            client.files.delete(name=uploaded_file.name)
        except:
            pass
        
        return response.text
        
    finally:
        # 一時ファイルの完全削除（例外発生時も実行）
        if os.path.exists(tmp_path):
            try:
                os.remove(tmp_path)
            except:
                pass
        
        # 分割チャンクファイルの削除
        for chunk in chunk_files:
            if os.path.exists(chunk):
                try:
                    os.remove(chunk)
                except:
                    pass

def generate_slides_from_minutes(minutes_content: str, model_name: str) -> dict:
    """議事録からスライド構造を生成"""
    client = genai.Client(api_key=st.secrets["api"]["google_api_key"])
    
    prompt = f"{SLIDE_PROMPT}\n\n---\n\n議事録:\n{minutes_content}"
    
    response = client.models.generate_content(
        model=model_name,
        contents=prompt
    )
    
    # JSONを抽出
    response_text = response.text
    try:
        # ```json ... ``` で囲まれている場合
        if "```json" in response_text:
            json_str = response_text.split("```json")[1].split("```")[0].strip()
        elif "```" in response_text:
            json_str = response_text.split("```")[1].split("```")[0].strip()
        else:
            json_str = response_text.strip()
        
        return json.loads(json_str)
    except Exception as e:
        st.error(f"スライド構造の解析エラー: {e}")
        return None

def convert_slides_to_markdown(slides_data: dict) -> str:
    """スライドデータをMarp形式Markdownに変換"""
    if not slides_data:
        return ""
    
    lines = [
        "---",
        "marp: true",
        "theme: default",
        "paginate: true",
        "---",
        ""
    ]
    
    for slide in slides_data.get("slides", []):
        slide_type = slide.get("type", "content")
        
        if slide_type == "title":
            lines.append(f"# {slide.get('title', '')}")
            if slide.get("subtitle"):
                lines.append(f"\n### {slide['subtitle']}")
        
        elif slide_type == "purpose":
            lines.append(f"## {slide.get('title', '会議の目的')}")
            lines.append(f"\n{slide.get('content', '')}")
        
        elif slide_type == "discussion":
            lines.append(f"## {slide.get('title', '')}")
            for bullet in slide.get("bullets", []):
                lines.append(f"- {bullet}")
        
        elif slide_type == "actions":
            lines.append(f"## {slide.get('title', 'ネクストアクション')}")
            for item in slide.get("items", []):
                assignee = item.get("assignee", "")
                task = item.get("task", "")
                deadline = item.get("deadline", "")
                deadline_str = f" ({deadline})" if deadline else ""
                lines.append(f"- **{assignee}**: {task}{deadline_str}")
        
        elif slide_type == "decisions":
            lines.append(f"## {slide.get('title', '決定事項')}")
            for bullet in slide.get("bullets", []):
                lines.append(f"- ✅ {bullet}")
        
        else:
            # 汎用コンテンツ
            lines.append(f"## {slide.get('title', '')}")
            for bullet in slide.get("bullets", []):
                lines.append(f"- {bullet}")
        
        lines.append("")
        lines.append("---")
        lines.append("")
    
    return "\n".join(lines)

def convert_slides_to_pptx(slides_data: dict) -> bytes:
    """スライドデータをPowerPoint形式に変換"""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        from io import BytesIO
    except ImportError:
        st.error("PowerPoint生成には python-pptx が必要です")
        return None
    
    if not slides_data:
        return None
    
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    for slide in slides_data.get("slides", []):
        slide_type = slide.get("type", "content")
        
        if slide_type == "title":
            layout = prs.slide_layouts[6]  # 空白
            ppt_slide = prs.slides.add_slide(layout)
            
            # タイトル
            title_box = ppt_slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = slide.get("title", "")
            p.font.size = Pt(44)
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER
            
            # サブタイトル
            if slide.get("subtitle"):
                sub_box = ppt_slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(0.8))
                tf = sub_box.text_frame
                p = tf.paragraphs[0]
                p.text = slide.get("subtitle", "")
                p.font.size = Pt(24)
                p.alignment = PP_ALIGN.CENTER
        
        else:
            layout = prs.slide_layouts[6]  # 空白
            ppt_slide = prs.slides.add_slide(layout)
            
            # タイトル
            title_box = ppt_slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(1))
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = slide.get("title", "")
            p.font.size = Pt(32)
            p.font.bold = True
            
            # コンテンツ
            content_box = ppt_slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.333), Inches(5.5))
            tf = content_box.text_frame
            tf.word_wrap = True
            
            if slide_type == "purpose":
                p = tf.paragraphs[0]
                p.text = slide.get("content", "")
                p.font.size = Pt(24)
            
            elif slide_type in ["discussion", "decisions"]:
                bullets = slide.get("bullets", [])
                for i, bullet in enumerate(bullets):
                    if i == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    prefix = "✅ " if slide_type == "decisions" else "• "
                    p.text = f"{prefix}{bullet}"
                    p.font.size = Pt(20)
                    p.space_after = Pt(12)
            
            elif slide_type == "actions":
                items = slide.get("items", [])
                for i, item in enumerate(items):
                    if i == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    assignee = item.get("assignee", "")
                    task = item.get("task", "")
                    deadline = item.get("deadline", "")
                    deadline_str = f" ({deadline})" if deadline else ""
                    p.text = f"• {assignee}: {task}{deadline_str}"
                    p.font.size = Pt(20)
                    p.space_after = Pt(12)
    
    # バイトストリームに保存
    output = BytesIO()
    prs.save(output)
    output.seek(0)
    return output.getvalue()

def generate_image_prompt(minutes_content: str, model_name: str) -> str:
    """議事録から画像生成用プロンプト（英語）を生成"""
    client = genai.Client(api_key=st.secrets["api"]["google_api_key"])
    
    prompt = f"""
    Based on the following meeting minutes, create a detailed prompt for an AI image generator (like Imagen 3) to create a visual summary slide.
    
    The image should be a professional infographics-style presentation slide.
    It should visually represent the key topics and decisions.
    
    Output ONLY the English prompt for the image generator. Do not include any other text.
    
    Minutes:
    {minutes_content[:1500]}
    """
    
    response = client.models.generate_content(
        model=model_name,
        contents=prompt
    )
    return response.text

def generate_slide_image(prompt: str) -> bytes:
    """Imagen 4.0で画像を生成（REST API使用）"""
    api_key = st.secrets["api"]["google_api_key"]
    # 利用可能な最新モデルに変更
    url = f"https://generativelanguage.googleapis.com/v1beta/models/imagen-4.0-generate-001:predict?key={api_key}"
    
    headers = {
        "Content-Type": "application/json"
    }
    
    data = {
        "instances": [
            {
                "prompt": prompt
            }
        ],
        "parameters": {
            "sampleCount": 1,
            "aspectRatio": "16:9"
        }
    }
    
    for attempt in range(3):
        try:
            # タイムアウトを設定してリクエスト
            response = requests.post(url, headers=headers, json=data, timeout=60)
            response_json = response.json()
            
            if "predictions" in response_json:
                # Base64エンコードされた画像データを取得
                b64_image = response_json["predictions"][0]["bytesBase64Encoded"]
                import base64
                return base64.b64decode(b64_image)
            else:
                error_msg = response_json.get("error", {}).get("message", "不明なエラー")
                if attempt < 2:
                    time.sleep(2)
                    continue
                st.error(f"画像生成APIエラー: {error_msg}")
                return None
                
        except Exception as e:
            if attempt < 2:
                time.sleep(2)
                continue
            st.error(f"画像生成リクエストエラー: {e}")
            return None

def generate_multi_slide_prompts(minutes_content: str, model_name: str) -> list:
    """議事録から複数枚のスライド用プロンプトを生成"""
    client = genai.Client(api_key=st.secrets["api"]["google_api_key"])
    
    # テキストモデル用プロンプト
    prompt = f"""
    Based on the following meeting minutes, create a plan for a presentation slide deck (approx. 3-6 slides).
    
    The output must be a JSON list of objects. Each object represents one slide and must have:
    - "title": Slide title in Japanese (approx 20 chars max).
    - "description": Bullet points or short description in Japanese (3-5 lines max).
    - "image_prompt": A prompt for Imagen 4.0 to generate a BACKGROUND image for this slide.
      IMPORTANT: The prompt must explicitly say "NO TEXT", "Minimalist abstract background", "Professional presentation background".
      Do NOT ask for any text in the image. The text will be added programmatically later.
      Focus on the visual style (e.g., "Corporate blue theme", "Clean white geometric patterns", "Soft gradient").
    
    Minutes:
    {minutes_content[:3000]}
    
    Output JSON only.
    """
    
    # リトライ設定
    max_retries = 3
    
    for attempt in range(max_retries):
        try:
            response = client.models.generate_content(
                model=model_name,
                contents=prompt,
                config=types.GenerateContentConfig(
                    response_mime_type="application/json"
                )
            )
            return json.loads(response.text)
        except Exception as e:
            if "RESOURCE_EXHAUSTED" in str(e) or "429" in str(e):
                if attempt < max_retries - 1:
                    wait_time = 60
                    st.warning(f"⚠️ API制限（構成案作成）。トークン回復まで{wait_time}秒待機します... ({attempt + 1}/{max_retries})")
                    time.sleep(wait_time)
                    continue
            
            # 最後の試行またはその他のエラー
            if attempt == max_retries - 1:
                st.error(f"スライド構成案の生成エラー: {e}")
                return []

def draw_text_on_image(image_bytes: bytes, title: str, description: str) -> bytes:
    """背景画像に日本語テキストを描画"""
    try:
        # 画像読み込み
        img = Image.open(io.BytesIO(image_bytes)).convert("RGBA")
        width, height = img.size
        
        # 半透明の黒オーバーレイを作成して文字を見やすくする
        overlay = Image.new("RGBA", img.size, (0, 0, 0, 0))
        draw = ImageDraw.Draw(overlay)
        
        # タイトル用帯（上部）
        draw.rectangle([(0, 0), (width, int(height * 0.2))], fill=(0, 0, 0, 100))
        # 本文用エリア（中央〜下）- 少し暗くする
        draw.rectangle([(0, int(height * 0.2)), (width, height)], fill=(0, 0, 0, 40))
        
        # 合成
        img = Image.alpha_composite(img, overlay)
        draw = ImageDraw.Draw(img)
        
        # フォント設定
        try:
            title_font = ImageFont.truetype("fonts/NotoSansJP-Bold.otf", int(height * 0.08))
            body_font = ImageFont.truetype("fonts/NotoSansJP-Regular.otf", int(height * 0.05))
        except OSError:
            # フォントがない場合はデフォルト（日本語が出ない可能性ありだがフォールバック）
            title_font = ImageFont.load_default()
            body_font = ImageFont.load_default()
        
        # タイトル描画（中央揃え）
        # textbbox (left, top, right, bottom)
        bbox = draw.textbbox((0, 0), title, font=title_font)
        text_w = bbox[2] - bbox[0]
        text_h = bbox[3] - bbox[1]
        
        title_x = (width - text_w) / 2
        title_y = (int(height * 0.2) - text_h) / 2
        draw.text((title_x, title_y), title, font=title_font, fill=(255, 255, 255, 255))
        
        # 本文描画（左寄せ、折り返しあり）
        # 簡易的な折り返し処理（文字数ベース）
        margin = int(width * 0.1)
        current_y = int(height * 0.3)
        chars_per_line = int((width - 2 * margin) / (int(height * 0.05)))  # およその文字数
        
        lines = description.split('\n')
        wrapped_lines = []
        for line in lines:
            if len(line) > chars_per_line:
                for j in range(0, len(line), chars_per_line):
                    wrapped_lines.append(line[j:j+chars_per_line])
            else:
                wrapped_lines.append(line)
        
        for line in wrapped_lines:
            draw.text((margin, current_y), line, font=body_font, fill=(255, 255, 255, 230))
            current_y += int(height * 0.08)  # 行間
        
        # バイト列に戻す
        output = io.BytesIO()
        img.convert("RGB").save(output, format="PNG")
        return output.getvalue()
        
    except Exception as e:
        st.warning(f"文字描画エラー: {e}")
        return image_bytes

def generate_slide_images_batch(slide_prompts: list) -> list:
    """複数のスライド画像をバッチ生成（文字合成付き）"""
    generated_images = []
    total = len(slide_prompts)
    progress_bar = st.progress(0)
    status_text = st.empty()
    
    for i, slide in enumerate(slide_prompts):
        title = slide.get("title", "Untitled")
        status_text.text(f"🖼️ 画像生成中 ({i+1}/{total}): {title}")
        img_prompt = slide.get("image_prompt", "")
        
        # 429対策の待機
        if i > 0:
            time.sleep(5)
            
        # 1. 背景画像を生成
        bg_bytes = generate_slide_image(img_prompt)
        
        if bg_bytes:
            # 2. 文字を合成（ハイブリッド処理）
            final_image = draw_text_on_image(bg_bytes, title, slide.get("description", ""))
            
            generated_images.append({
                "title": title,
                "image": final_image
            })
        else:
            st.warning(f"スライド「{title}」の生成に失敗しました（スキップします）")
        
        progress_bar.progress((i + 1) / total)
        
    status_text.empty()
    progress_bar.empty()
    return generated_images

def create_slides_zip(generated_images: list) -> bytes:
    """画像をZIPにまとめる"""
    zip_buffer = io.BytesIO()
    with zipfile.ZipFile(zip_buffer, "w", zipfile.ZIP_DEFLATED) as zf:
        for i, item in enumerate(generated_images):
            # ファイル名: 01_Title.png
            safe_title = "".join(c for c in item["title"] if c.isalnum() or c in (' ', '_', '-')).strip()
            filename = f"{i+1:02d}_{safe_title}.png"
            zf.writestr(filename, item["image"])
    return zip_buffer.getvalue()

def show_slide_generator(minutes_content: str, model_name: str):
    """スライド生成UIを表示（画像スライドのみ）"""
    st.markdown("---")
    st.markdown("### 📊 スライド生成")
    
    st.markdown("議事録の内容に基づき、**複数のスライド（画像）を自動生成**します。タイトルや説明文も読みやすく配置されます。")

    if st.button("📊 スライドを生成 (Imagen 4.0)", type="primary", use_container_width=True):
        # 複数枚生成モード
        with st.spinner("スライド構成案を作成中..."):
            slide_prompts = generate_multi_slide_prompts(minutes_content, model_name)
        
        if slide_prompts:
            st.info(f"💡 {len(slide_prompts)}枚のスライド構成案を作成しました。画像を生成します...")
            images = generate_slide_images_batch(slide_prompts)
            
            if images:
                st.success(f"✨ {len(images)}枚のスライド画像を生成しました！")
                
                # ギャラリー表示
                for img in images:
                    st.image(img["image"], caption=img["title"], use_container_width=True)
                
                # ZIPダウンロード
                zip_data = create_slides_zip(images)
                st.download_button(
                    label="📥 全画像をダウンロード (.zip)",
                    data=zip_data,
                    file_name="slides_images.zip",
                    mime="application/zip",
                    use_container_width=True
                )
            else:
                st.error("画像の生成に失敗しました。")
        else:
            st.error("スライド構成案の作成に失敗しました。")

def main():
    st.title("📊 MinuteSlide")
    st.markdown("**議事録からスライドへ、一瞬で変換。** 音声や動画をアップロードするだけで、AIが詳細な議事録とプレゼン資料を自動生成します。")
    
    # Firestore初期化
    db = init_firestore()
    
    # Firebase Authentication（組織対応）
    if not check_firebase_auth(db):
        return
    
    # 組織情報確認
    org_id = st.session_state.organization_id
    org_name = st.session_state.organization_name
    
    # 組織未登録の場合は登録画面を表示
    if org_id is None and db:
        st.warning("⚠️ 組織情報が登録されていません")
        st.markdown("### 組織の登録")
        st.markdown("議事録を保存・共有するには組織（会社名・チーム名）の登録が必要です。")
        
        new_org_name = st.text_input("組織名（会社名・チーム名）", 
                                      placeholder="例：株式会社ABC / 営業チーム",
                                      help="同じ組織名のメンバーと議事録を共有できます")
        
        col1, col2 = st.columns(2)
        with col1:
            if st.button("組織を登録", type="primary", use_container_width=True):
                if new_org_name and len(new_org_name.strip()) >= 2:
                    with st.spinner("登録中..."):
                        org_id = get_or_create_organization(db, new_org_name)
                        save_user_organization(
                            db, 
                            st.session_state.user_id, 
                            st.session_state.user_email, 
                            org_id, 
                            new_org_name.strip()
                        )
                        st.session_state.organization_id = org_id
                        st.session_state.organization_name = new_org_name.strip()
                    st.success("組織を登録しました！")
                    st.rerun()
                else:
                    st.error("組織名を2文字以上で入力してください")
        with col2:
            if st.button("ログアウト", use_container_width=True):
                st.session_state.authenticated = False
                st.session_state.user_email = None
                st.session_state.user_id = None
                st.session_state.organization_id = None
                st.session_state.organization_name = None
                st.rerun()
        return
    
    # サイドバー
    st.sidebar.header("⚙️ 設定")
    st.sidebar.success(f"✅ {st.session_state.user_email}")
    if org_name:
        st.sidebar.info(f"🏢 {org_name}")
    
    # ログアウトボタン
    if st.sidebar.button("ログアウト"):
        st.session_state.authenticated = False
        st.session_state.user_email = None
        st.session_state.user_id = None
        st.session_state.organization_id = None
        st.session_state.organization_name = None
        st.rerun()
    
    st.sidebar.markdown("---")
    
    # モデル選択
    model_option = st.sidebar.selectbox(
        "モデル選択",
        options=["gemini-2.0-flash", "gemini-2.0-flash-lite-001", "gemini-2.5-flash"],
        format_func=lambda x: f"{x} (軽量)" if "lite" in x else f"{x} (最新)"
    )
    
    st.sidebar.markdown("---")
    st.sidebar.markdown("""
    **対応形式:** .mp3, .wav, .m4a  
    **最大サイズ:** 500MB
    """)
    
    # 履歴表示（同一組織のみ）
    if db and org_id:
        st.sidebar.markdown("---")
        st.sidebar.subheader("📚 議事録履歴")
        history = get_minutes_history(db, org_id)
        
        if history:
            for item in history:
                created_at = item.get("created_at")
                if created_at:
                    date_str = created_at.strftime("%m/%d %H:%M")
                else:
                    date_str = "日時不明"
                
                label = f"{date_str} - {item.get('title', '無題')[:15]}"
                if st.sidebar.button(label, key=item["id"]):
                    st.session_state.selected_minute = item["id"]
                    st.rerun()
        else:
            st.sidebar.caption("履歴がありません")
    
    # メインエリア
    if "selected_minute" in st.session_state and st.session_state.selected_minute and db and org_id:
        minute = get_minute_by_id(db, org_id, st.session_state.selected_minute)
        if minute:
            st.markdown("## 📋 過去の議事録")
            col1, col2 = st.columns([3, 1])
            with col1:
                st.markdown(f"**タイトル:** {minute.get('title', '無題')}")
                st.markdown(f"**作成者:** {minute.get('created_by', '不明')}")
                if minute.get("created_at"):
                    st.markdown(f"**作成日時:** {minute['created_at'].strftime('%Y-%m-%d %H:%M')}")
            with col2:
                if st.button("🆕 新規作成に戻る"):
                    st.session_state.selected_minute = None
                    st.rerun()
                
                # 削除ボタン
                if st.button("🗑️ 削除", type="primary"):
                    if delete_minute(db, org_id, st.session_state.selected_minute):
                        st.success("削除しました")
                        st.session_state.selected_minute = None
                        time.sleep(1)
                        st.rerun()
                    else:
                        st.error("削除に失敗しました")
            
            st.markdown("---")
            st.markdown(minute.get("content", ""))
            
            # スライド生成UI
            show_slide_generator(minute.get("content", ""), model_option)
            
            st.markdown("---")
            st.code(minute.get("content", ""), language="markdown")
            return
    
    # 新規議事録作成
    col1, col2 = st.columns([2, 1])
    
    with col1:
        uploaded_file = st.file_uploader(
            "音声ファイルをアップロード",
            type=["mp3", "wav", "m4a", "aac", "ogg", "flac"],
            help="会議の音声ファイルをドラッグ＆ドロップしてください"
        )
        
        if uploaded_file:
            st.audio(uploaded_file)
            st.caption(f"📁 {uploaded_file.name} ({uploaded_file.size / 1024 / 1024:.2f} MB)")
    
    with col2:
        title = st.text_input(
            "議事録タイトル",
            placeholder="例：週次定例会議 2024/01/25"
        )
        
        additional_instructions = st.text_area(
            "追加指示（任意）",
            placeholder="例：専門用語が多いので、技術的な文脈を重視して",
            height=80
        )
    
    # 実行ボタン
    if st.button("🚀 音声を解析して議事録作成", type="primary", use_container_width=True):
        if not uploaded_file:
            st.error("音声ファイルをアップロードしてください")
            return
        
        if not org_id:
            st.error("組織情報がありません。ログアウトして再度登録してください。")
            return
        
        with st.spinner("議事録を生成中..."):
            result = upload_and_process_audio(
                uploaded_file, 
                model_option, 
                additional_instructions
            )
        
        if result:
            st.success("議事録が生成されました！")
            
            # Firestoreに保存
            if db:
                doc_title = title if title else uploaded_file.name
                doc_id = save_to_firestore(
                    db, 
                    doc_title, 
                    result, 
                    uploaded_file.name,
                    st.session_state.user_email,
                    st.session_state.user_id,
                    org_id
                )
                if doc_id:
                    st.success("議事録を保存しました！詳細画面へ移動します...")
                    st.session_state.selected_minute = doc_id
                    time.sleep(1)
                    st.rerun()
            
            # 結果表示
            st.markdown("---")
            st.markdown("## 📋 生成された議事録")
            st.markdown(result)
            
            st.markdown("---")
            st.markdown("### 📝 コピー用（Markdown形式）")
            st.code(result, language="markdown")
            
            st.download_button(
                label="📥 議事録をダウンロード (.md)",
                data=result,
                file_name="meeting_minutes.md",
                mime="text/markdown"
            )
            
            # スライド生成UI
            show_slide_generator(result, model_option)

if __name__ == "__main__":
    main()
